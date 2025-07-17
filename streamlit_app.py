import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
import base64
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
import math
from datetime import datetime

st.set_page_config(layout="wide")
st.title("Advanced Current and Power Analysis in Wind Turbines")

# --- File Upload ---
st.header("1. Upload Data Files")
current_file = st.file_uploader("Upload Current Data File (XLSX)", type=["xlsx"])
voltage_file = st.file_uploader("Upload Voltage Data File (XLSX)", type=["xlsx"])
production_file = st.file_uploader("Upload Production Data File (XLSX)", type=["xlsx"])

# --- Analysis Settings ---
st.sidebar.header("2. Analysis Settings")
watermark_text = st.sidebar.text_input("Watermark Text", "Confidential Analysis")
watermark_opacity = st.sidebar.slider("Watermark Opacity", 0.1, 1.0, 0.2)

st.sidebar.subheader("Overload and Imbalance Thresholds")
num_turbines = st.sidebar.number_input("Number of Turbines", min_value=1, max_value=10, value=4)
default_overload_threshold = 2630

st.sidebar.write("Configure Overload Thresholds (Amperes)")
user_overload_thresholds = {}
for i in range(1, num_turbines + 1):
    turbine_id = f"WTG{i:02d}"
    user_overload_thresholds[turbine_id] = st.sidebar.number_input(
        f"Overload Threshold {turbine_id} (A)",
        min_value=1.0,
        value=float(default_overload_threshold),
        step=10.0,
        help=f"Enter the maximum allowable current for {turbine_id}."
    )
overload_thresholds = user_overload_thresholds

min_overload_duration_minutes = st.sidebar.slider(
    "Minimum Sustained Overload Duration (minutes)", 0, 60, 5
)

imbalance_threshold = st.sidebar.slider("Imbalance Threshold (%)", 1, 50, 10, 1)
efficiency_threshold = st.sidebar.slider("Efficiency Threshold (%)", 50, 100, 85, 1)

# New feature: Date range selector
st.sidebar.subheader("Date Range Filter")
date_range = st.sidebar.date_input(
    "Select Date Range",
    value=[datetime.now().replace(day=1), datetime.now()],
    min_value=datetime(2020, 1, 1),
    max_value=datetime.now()
)

if current_file and voltage_file and production_file:
    # --- Data Loading and Preprocessing ---
    st.header("2. Data Loading and Preprocessing")
    
    def load_and_preprocess_data(file, rename_dict, data_type):
        try:
            df = pd.read_excel(file)
            df.rename(columns=rename_dict, inplace=True)
            
            if 'PCTimeStamp' in df.columns:
                df['Date_Time'] = pd.to_datetime(
                    df['PCTimeStamp'],
                    format="%d.%m.%Y %H:%M:%S",
                    errors='coerce'
                )
                df['Date_Time'] = df['Date_Time'].fillna(
                    pd.to_datetime(df['PCTimeStamp'], errors='coerce'))
                df = df.dropna(subset=['Date_Time'])
                
                if df['Date_Time'].isnull().any():
                    st.warning(f"Some 'PCTimeStamp' values in {data_type} data could not be parsed.")
            else:
                st.error(f"{data_type} data missing 'PCTimeStamp' column.")
                return None
            
            # Apply date range filter
            if len(date_range) == 2:
                # Ensure date_range elements are datetime objects for proper comparison
                start_date = pd.to_datetime(date_range[0])
                end_date = pd.to_datetime(date_range[1])
                mask = (df['Date_Time'] >= start_date) & \
                       (df['Date_Time'] <= end_date)
                df = df.loc[mask]
            
            st.success(f"{data_type} data loaded successfully.")
            return df
            
        except Exception as e:
            st.error(f"Error loading {data_type} data: {e}")
            return None

    # Current data
    current_rename_dict = {
        'WTG01_Grid  CurrentPhase1 Max. (1)': 'WTG01A',
        'WTG01_Grid  CurrentPhase2 Max. (5)': 'WTG01B',
        'WTG01_Grid  CurrentPhase3 Max. (9)': 'WTG01C',
        'WTG02_Grid  CurrentPhase1 Max. (2)': 'WTG02A',
        'WTG02_Grid  CurrentPhase2 Max. (6)': 'WTG02B',
        'WTG02_Grid  CurrentPhase3 Max. (10)': 'WTG02C',
        'WTG03_Grid  CurrentPhase1 Max. (3)': 'WTG03A',
        'WTG03_Grid  CurrentPhase2 Max. (7)': 'WTG03B',
        'WTG03_Grid  CurrentPhase3 Max. (11)': 'WTG03C',
        'WTG04_Grid  CurrentPhase1 Max. (4)': 'WTG04A',
        'WTG04_Grid  CurrentPhase2 Max. (8)': 'WTG04B',
        'WTG04_Grid CurrentPhase3 Max. (12)': 'WTG04C'
    }
    current_df = load_and_preprocess_data(current_file, current_rename_dict, "Current")
    
    # Voltage data
    voltage_rename_dict = {
        'WTG01_Grid Production VoltagePhase1 Min. (1)': 'WTG01Vab',
        'WTG01_Grid Production VoltagePhase2 Min. (5)': 'WTG01Vbc',
        'WTG01_Grid Production VoltagePhase3 Min. (9)': 'WTG01Vca',
        'WTG02_Grid Production VoltagePhase1 Min. (2)': 'WTG02Vab',
        'WTG02_Grid Production VoltagePhase2 Min. (6)': 'WTG02Vbc',
        'WTG02_Grid Production VoltagePhase3 Min. (10)': 'WTG02Vca',
        'WTG03_Grid Production VoltagePhase1 Min. (3)': 'WTG03Vab',
        'WTG03_Grid Production VoltagePhase2 Min. (7)': 'WTG03Vbc',
        'WTG03_Grid Production VoltagePhase3 Min. (11)': 'WTG03Vca',
        'WTG04_Grid Production VoltagePhase1 Min. (4)': 'WTG04Vab',
        'WTG04_Grid Production VoltagePhase2 Min. (8)': 'WTG04Vbc',
        'WTG04_Grid Production VoltagePhase3 Min. (12)': 'WTG04Vca'
    }
    voltage_df = load_and_preprocess_data(voltage_file, voltage_rename_dict, "Voltage")
    
    # Production data
    production_rename_dict = {
        'WTG01_Production LatestAverage Total Active Power Avg. (1)': 'WTG01P',
        'WTG02_Production LatestAverage Total Active Power Avg. (2)': 'WTG02P',
        'WTG03_Production LatestAverage Total Active Power Avg. (3)': 'WTG03P',
        'WTG04_Production LatestAverage Total Active Power Avg. (4)': 'WTG04P'
    }
    production_df = load_and_preprocess_data(production_file, production_rename_dict, "Production")
    
    if production_df is not None:
        for col in ['WTG01P', 'WTG02P', 'WTG03P', 'WTG04P']:
            if col in production_df.columns:
                production_df[col] = production_df[col].clip(lower=0)

    if current_df is not None and voltage_df is not None and production_df is not None:
        try:
            df = pd.merge(current_df, voltage_df, on='Date_Time', how='outer')
            df = pd.merge(df, production_df, on='Date_Time', how='outer')
            df = df.sort_values('Date_Time').reset_index(drop=True)
            df.drop_duplicates(subset=['Date_Time'], inplace=True)
            df.set_index('Date_Time', inplace=True)
            df = df[~df.index.isna()]
            
            st.success("Dataframes merged successfully.")
            st.write("First 5 rows of merged data:")
            st.dataframe(df.head())

        except Exception as e:
            st.error(f"Error during data merging: {e}")
            st.stop()

        # --- Data Analysis ---
        st.header("3. Data Analysis")

        results = {}
        turbine_ids = [f"WTG{i:02d}" for i in range(1, num_turbines + 1)]

        numeric_cols = []
        for turbine_id in turbine_ids:
            numeric_cols.extend([f'{turbine_id}A', f'{turbine_id}B', f'{turbine_id}C'])
            numeric_cols.extend([f'{turbine_id}Vab', f'{turbine_id}Vbc', f'{turbine_id}Vca'])
            numeric_cols.append(f'{turbine_id}P')
            
        for col in numeric_cols:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce')
                df[col] = df[col].fillna(0)

        for turbine_id in turbine_ids:
            col_R = f'{turbine_id}A'
            col_S = f'{turbine_id}B'
            col_T = f'{turbine_id}C'
            col_Vab = f'{turbine_id}Vab'
            col_Vbc = f'{turbine_id}Vbc'
            col_Vca = f'{turbine_id}Vca'
            col_Prod = f'{turbine_id}P'

            # Calculate total current (vector sum divided by sqrt(3))
            df[f'{turbine_id}_Total'] = np.sqrt(df[col_R]**2 + df[col_S]**2 + df[col_T]**2) / np.sqrt(3)

            # Current imbalance calculations
            current_phases = df[[col_R, col_S, col_T]]
            df[f'{turbine_id}_RvgCurrent'] = current_phases.mean(axis=1)
            df[f'{turbine_id}_MaxDeviation'] = current_phases.max(axis=1) - df[f'{turbine_id}_RvgCurrent']
            df[f'{turbine_id}_Imbalance'] = (df[f'{turbine_id}_MaxDeviation'] / df[f'{turbine_id}_RvgCurrent']) * 100
            df[f'{turbine_id}_Imbalance'].replace([np.inf, -np.inf], np.nan, inplace=True)
            df[f'{turbine_id}_Imbalance'].fillna(0, inplace=True)

            # Voltage and power calculations
            if all(col in df.columns for col in [col_Vab, col_Vbc, col_Vca]):
                df[f'{turbine_id}_RvgVoltage'] = df[[col_Vab, col_Vbc, col_Vca]].mean(axis=1)
                df[f'{turbine_id}_RpparentPower_VA'] = np.sqrt(3) * df[f'{turbine_id}_RvgVoltage'] * df[f'{turbine_id}_Total']
                df[f'{turbine_id}_RpparentPower_kVA'] = df[f'{turbine_id}_RpparentPower_VA'] / 1000
            else:
                df[f'{turbine_id}_RpparentPower_kVA'] = 0

            # Power factor calculation
            if col_Prod in df.columns and f'{turbine_id}_RpparentPower_kVA' in df.columns:
                df[f'{turbine_id}_PF'] = df[col_Prod] / df[f'{turbine_id}_RpparentPower_kVA'].replace(0, np.nan)
                df[f'{turbine_id}_PF'].replace([np.inf, -np.inf], np.nan, inplace=True)
                df[f'{turbine_id}_PF'].fillna(0, inplace=True)
                df[f'{turbine_id}_PF'] = df[f'{turbine_id}_PF'].clip(-1.0, 1.0)
            else:
                df[f'{turbine_id}_PF'] = 0

            # Energy production estimation
            if col_Prod in df.columns:
                # Assuming 10-minute intervals, so (10/60) hours for each reading
                df[f'{turbine_id}_Daily_MWh'] = df[col_Prod] * (10/60) / 1000
                total_mwh_period = df[f'{turbine_id}_Daily_MWh'].sum()
                st.write(f"Estimated total MWh for {turbine_id} over the data period: {total_mwh_period:,.2f} MWh")

            # Overload detection
            df[f'{turbine_id}_Overload_R'] = df[col_R] > overload_thresholds[turbine_id]
            df[f'{turbine_id}_Overload_S'] = df[col_S] > overload_thresholds[turbine_id]
            df[f'{turbine_id}_Overload_T'] = df[col_T] > overload_thresholds[turbine_id]
            df[f'{turbine_id}_Overload_Total'] = df[f'{turbine_id}_Total'] > overload_thresholds[turbine_id]
            df[f'{turbine_id}_Rny_Overload'] = df[f'{turbine_id}_Overload_R'] | df[f'{turbine_id}_Overload_S'] | \
                                               df[f'{turbine_id}_Overload_T'] | df[f'{turbine_id}_Overload_Total']

            # Sustained overload detection
            if len(df.index) > 1:
                interval_seconds = (df.index[1] - df.index[0]).total_seconds()
            else:
                interval_seconds = 600 # Default to 10 minutes if only one data point

            window_size_points = max(1, int((min_overload_duration_minutes * 60) / interval_seconds))
            df[f'{turbine_id}_Sustained_Overload'] = df[f'{turbine_id}_Rny_Overload'].rolling(
                window=window_size_points, min_periods=window_size_points
            ).sum() >= window_size_points

            # Imbalance and efficiency events
            df[f'{turbine_id}_Imbalance_Event'] = df[f'{turbine_id}_Imbalance'] > imbalance_threshold
            df[f'{turbine_id}_Efficiency_Event'] = df[f'{turbine_id}_PF'] < (efficiency_threshold / 100)

            # Store results
            results[turbine_id] = {
                "Overload_Events": df[df[f'{turbine_id}_Sustained_Overload']].index.tolist(),
                "Imbalance_Events": df[df[f'{turbine_id}_Imbalance_Event']].index.tolist(),
                "Efficiency_Events": df[df[f'{turbine_id}_Efficiency_Event']].index.tolist(),
                "Peak_Current_R": df[col_R].max(),
                "Peak_Current_S": df[col_S].max(),
                "Peak_Current_T": df[col_T].max(),
                "Peak_Total_Current": df[f'{turbine_id}_Total'].max(),
                "Avg_Imbalance": df[f'{turbine_id}_Imbalance'].mean(),
                "Avg_Power_Factor": df[f'{turbine_id}_PF'].mean(),
                "Annual_Mwh_Estimate": total_mwh_period # This is actually total_mwh_period, not annual
            }

        st.success("Analysis complete for all available turbines.")

        # Display Summary Results for Streamlit (unchanged from previous request)
        st.header("4. Summary of Analysis Results")
        for turbine_id, data in results.items():
            st.subheader(f"Turbine {turbine_id} Summary:")
            st.write(f"- Peak Current (Phase R): {data['Peak_Current_R']:.2f} A")
            st.write(f"- Peak Current (Phase S): {data['Peak_Current_S']:.2f} A")
            st.write(f"- Peak Current (Phase T): {data['Peak_Current_T']:.2f} A")
            st.write(f"- Peak Total Current: {data['Peak_Total_Current']:.2f} A (Threshold: {overload_thresholds.get(turbine_id, 'N/A')} A)")
            st.write(f"- Average Current Imbalance: {data['Avg_Imbalance']:.2f}% (Threshold: {imbalance_threshold}%)")
            st.write(f"- Average Power Factor: {data['Avg_Power_Factor']:.2f} (Threshold: {efficiency_threshold}%)")
            st.write(f"- Sustained Overload Events: {len(data['Overload_Events'])} (duration >= {min_overload_duration_minutes} min)")
            st.write(f"- Imbalance Events: {len(data['Imbalance_Events'])}")
            st.write(f"- Low Efficiency Events: {len(data['Efficiency_Events'])}")
            st.write(f"- Estimated Total MWh over data period: {data['Annual_Mwh_Estimate']:.2f} MWh")

            # Display selected events in Streamlit table
            event_data = []
            for event_type, events in [
                ("Sustained Overload", data['Overload_Events']),
                ("High Imbalance", data['Imbalance_Events']) # Only these two
            ]:
                for ts in events:
                    if pd.isna(ts):
                        continue
                    row = df.loc[ts]
                    event_detail = {
                        "Date_Time": ts.strftime('%Y-%m-%d %H:%M') if not pd.isna(ts) else "Invalid Date",
                        "Event_Type": event_type,
                        "Phase_R_Current": f"{row[f'{turbine_id}A']:.2f} A" if f'{turbine_id}A' in row else "N/A",
                        "Phase_S_Current": f"{row[f'{turbine_id}B']:.2f} A" if f'{turbine_id}B' in row else "N/A",
                        "Phase_T_Current": f"{row[f'{turbine_id}C']:.2f} A" if f'{turbine_id}C' in row else "N/A",
                        "Total_Current": f"{row[f'{turbine_id}_Total']:.2f} A" if f'{turbine_id}_Total' in row else "N/A",
                        "Imbalance": f"{row[f'{turbine_id}_Imbalance']:.2f}%" if f'{turbine_id}_Imbalance' in row else "N/A",
                        "Power_Factor": f"{row[f'{turbine_id}_PF']:.2f}" if f'{turbine_id}_PF' in row else "N/A"
                    }
                    event_data.append(event_detail)

            if event_data:
                st.dataframe(pd.DataFrame(event_data))
            else:
                st.success("No significant 'High Imbalance' or 'Sustained Overload' events detected in this turbine based on current thresholds.")


        # --- Visualization ---
        st.header("5. Data Visualization")

        def plot_currents_voltages(df_plot, turbine_id, overload_threshold):
            fig, ax1 = plt.subplots(figsize=(15, 7))

            ax1.set_xlabel('Date Time')
            ax1.set_ylabel('Current (A)', color='tab:blue')
            ax1.plot(df_plot.index, df_plot[f'{turbine_id}A'], label=f'{turbine_id} Current R', color='blue', alpha=0.7)
            ax1.plot(df_plot.index, df_plot[f'{turbine_id}B'], label=f'{turbine_id} Current S', color='green', alpha=0.7)
            ax1.plot(df_plot.index, df_plot[f'{turbine_id}C'], label=f'{turbine_id} Current T', color='orange', alpha=0.7)
            ax1.plot(df_plot.index, df_plot[f'{turbine_id}_Total'], label=f'{turbine_id} Total Current', color='purple', linestyle='--', linewidth=2)
            ax1.axhline(y=overload_threshold, color='red', linestyle=':', linewidth=2, label=f'Overload Threshold ({overload_threshold}A)')
            ax1.tick_params(axis='y', labelcolor='tab:blue')

            if f'{turbine_id}_Sustained_Overload' in df_plot.columns:
                overload_points = df_plot[df_plot[f'{turbine_id}_Sustained_Overload']]
                if not overload_points.empty:
                    ax1.scatter(overload_points.index, overload_points[f'{turbine_id}_Total'], color='red', s=50, zorder=5, label='Sustained Overload Event')

            ax2 = ax1.twinx()
            ax2.set_ylabel('Voltage (V)', color='tab:red')
            if f'{turbine_id}Vab' in df_plot.columns:
                ax2.plot(df_plot.index, df_plot[f'{turbine_id}Vab'], label=f'{turbine_id} Voltage Vab', color='red', alpha=0.5)
                ax2.plot(df_plot.index, df_plot[f'{turbine_id}Vbc'], label=f'{turbine_id} Voltage Vbc', color='brown', alpha=0.5)
                ax2.plot(df_plot.index, df_plot[f'{turbine_id}Vca'], label=f'{turbine_id} Voltage Vca', color='pink', alpha=0.5)
            ax2.tick_params(axis='y', labelcolor='tab:red')

            ax3 = ax1.twinx()
            ax3.spines['right'].set_position(('outward', 60))
            ax3.set_ylabel('Power Factor', color='tab:green')
            if f'{turbine_id}_PF' in df_plot.columns:
                ax3.plot(df_plot.index, df_plot[f'{turbine_id}_PF'], label=f'{turbine_id} Power Factor', color='green', linestyle='-.', alpha=0.7)
            ax3.tick_params(axis='y', labelcolor='tab:green')

            fig.suptitle(f'{turbine_id} Current, Voltage, and Power Factor Over Time')
            fig.legend(loc="upper left", bbox_to_anchor=(0.05,0.95))
            fig.autofmt_xdate()
            plt.grid(True)
            st.pyplot(fig)

        def plot_imbalance(df_plot, turbine_id, imbalance_threshold):
            fig, ax = plt.subplots(figsize=(15, 7))
            ax.plot(df_plot.index, df_plot[f'{turbine_id}_Imbalance'], label=f'{turbine_id} Current Imbalance', color='teal')
            ax.axhline(y=imbalance_threshold, color='red', linestyle=':', label=f'Imbalance Threshold ({imbalance_threshold}%)')
            ax.set_xlabel('Date Time')
            ax.set_ylabel('Imbalance (%)')
            ax.set_title(f'{turbine_id} Current Imbalance Over Time')
            ax.legend()
            ax.grid(True)

            if f'{turbine_id}_Imbalance_Event' in df_plot.columns:
                imbalance_points = df_plot[df_plot[f'{turbine_id}_Imbalance_Event']]
                if not imbalance_points.empty:
                    ax.scatter(imbalance_points.index, imbalance_points[f'{turbine_id}_Imbalance'], color='red', s=50, zorder=5, label='Imbalance Event')

            fig.autofmt_xdate()
            st.pyplot(fig)

        for turbine_id in turbine_ids:
            if turbine_id in results:
                st.subheader(f"Graphs for {turbine_id}")
                plot_currents_voltages(df, turbine_id, overload_thresholds[turbine_id])
                plot_imbalance(df, turbine_id, imbalance_threshold)

        # --- Report Generation Functions ---
        def fig_to_bytes(fig):
            buf = BytesIO()
            fig.savefig(buf, format="png", bbox_inches="tight", dpi=100)
            buf.seek(0)
            return buf.getvalue()

        def create_current_pdf_report(analysis_results, df_full, watermark_text, watermark_opacity):
            pdf = FPDF()
            pdf.add_page()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.set_font("Arial", size=12)

            if watermark_text:
                pdf.set_font("Arial", "B", 50)
                gray_value = int(255 * (1 - watermark_opacity)) 
                pdf.set_text_color(gray_value, gray_value, gray_value)
                pdf.rotate(45, pdf.w / 2, pdf.h / 2)
                pdf.text(pdf.w / 4, pdf.h / 2, watermark_text)
                pdf.rotate(0)
                pdf.set_text_color(0, 0, 0) # Reset text color to black for content

            pdf.set_font("Arial", "B", 16)
            pdf.cell(0, 10, "Advanced Current and Power Analysis Report", ln=True, align="C")
            pdf.ln(10)

            pdf.set_font("Arial", "B", 14)
            pdf.cell(0, 10, "1. Summary of Results", ln=True)
            pdf.set_font("Arial", "", 10)

            for turbine_id, data in analysis_results.items():
                pdf.set_font("Arial", "B", 12)
                pdf.cell(0, 10, f"Turbine {turbine_id} Summary:", ln=True)
                pdf.set_font("Arial", "", 10)
                pdf.multi_cell(0, 6, f"- Peak Current (Phase R): {data['Peak_Current_R']:.2f} A")
                pdf.multi_cell(0, 6, f"- Peak Current (Phase S): {data['Peak_Current_S']:.2f} A")
                pdf.multi_cell(0, 6, f"- Peak Current (Phase T): {data['Peak_Current_T']:.2f} A")
                pdf.multi_cell(0, 6, f"- Peak Total Current: {data['Peak_Total_Current']:.2f} A")
                pdf.multi_cell(0, 6, f"- Average Current Imbalance: {data['Avg_Imbalance']:.2f}%")
                pdf.multi_cell(0, 6, f"- Average Power Factor: {data['Avg_Power_Factor']:.2f}")
                pdf.multi_cell(0, 6, f"- Sustained Overload Events: {len(data['Overload_Events'])}")
                pdf.multi_cell(0, 6, f"- Imbalance Events: {len(data['Imbalance_Events'])}")
                pdf.multi_cell(0, 6, f"- Low Efficiency Events: {len(data['Efficiency_Events'])}")
                pdf.multi_cell(0, 6, f"- Estimated Total MWh over data period: {data['Annual_Mwh_Estimate']:.2f} MWh")
                pdf.ln(5)

                # --- ONLY High Imbalance and Overload Events for PDF ---
                if data['Overload_Events'] or data['Imbalance_Events']:
                    pdf.set_font("Arial", "B", 10)
                    pdf.cell(0, 7, f"Detailed Overload & High Imbalance Events for {turbine_id}:", ln=True)
                    pdf.set_font("Arial", "", 8)
                    event_headers = ["Date_Time", "Event_Type", "Phase_R_Current", "Phase_S_Current", "Phase_T_Current", "Total_Current", "Imbalance", "Power_Factor"]
                    col_widths = [25, 25, 20, 20, 20, 20, 20, 20]

                    for i, header in enumerate(event_headers):
                        pdf.cell(col_widths[i], 7, header, 1, 0, 'C')
                    pdf.ln()

                    event_data_rows = []
                    # Include both Overload_Events and Imbalance_Events
                    for event_type, events in [
                        ("Sustained Overload", data['Overload_Events']),
                        ("High Imbalance", data['Imbalance_Events'])
                    ]:
                        for ts in events:
                            if pd.isna(ts):
                                continue
                            row_data = df_full.loc[ts]
                            event_row = [
                                ts.strftime('%Y-%m-%d %H:%M') if not pd.isna(ts) else "Invalid Date",
                                event_type,
                                f"{row_data.get(f'{turbine_id}A', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}B', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}C', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}_Total', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}_Imbalance', 0.0):.2f}%",
                                f"{row_data.get(f'{turbine_id}_PF', 0.0):.2f}"
                            ]
                            event_data_rows.append(event_row)
                    
                    # Sort events by Date_Time for chronological order in the report
                    event_data_rows.sort(key=lambda x: datetime.strptime(x[0], '%Y-%m-%d %H:%M'))

                    for row in event_data_rows:
                        for i, item in enumerate(row):
                            pdf.cell(col_widths[i], 6, str(item), 1, 0, 'C')
                        pdf.ln()
                    pdf.ln(5)

            pdf.add_page()
            if watermark_text:
                pdf.set_font("Arial", "B", 50)
                gray_value = int(255 * (1 - watermark_opacity))
                pdf.set_text_color(gray_value, gray_value, gray_value)
                pdf.rotate(45, pdf.w / 2, pdf.h / 2)
                pdf.text(pdf.w / 4, pdf.h / 2, watermark_text)
                pdf.rotate(0)
                pdf.set_text_color(0, 0, 0)

            pdf.set_font("Arial", "B", 14)
            pdf.cell(0, 10, "2. Data Visualizations", ln=True)
            pdf.ln(5)

            for turbine_id in analysis_results.keys():
                pdf.set_font("Arial", "B", 12)
                pdf.cell(0, 10, f"Graphs for Turbine {turbine_id}:", ln=True)

                fig_current_voltage, ax_cv = plt.subplots(figsize=(15, 7))
                ax_cv.set_xlabel('Date Time')
                ax_cv.set_ylabel('Current (A)', color='tab:blue')
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}A'], label=f'{turbine_id} Current R', color='blue', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}B'], label=f'{turbine_id} Current S', color='green', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}C'], label=f'{turbine_id} Current T', color='orange', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}_Total'], label=f'{turbine_id} Total Current', color='purple', linestyle='--', linewidth=2)
                ax_cv.axhline(y=overload_thresholds[turbine_id], color='red', linestyle=':', linewidth=2, label=f'Overload Threshold ({overload_thresholds[turbine_id]}A)')
                ax_cv.tick_params(axis='y', labelcolor='tab:blue')

                ax_cv2 = ax_cv.twinx()
                ax_cv2.set_ylabel('Voltage (V)', color='tab:red')
                if f'{turbine_id}Vab' in df_full.columns:
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vab'], label=f'{turbine_id} Voltage Vab', color='red', alpha=0.5)
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vbc'], label=f'{turbine_id} Voltage Vbc', color='brown', alpha=0.5)
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vca'], label=f'{turbine_id} Voltage Vca', color='pink', alpha=0.5)
                ax_cv2.tick_params(axis='y', labelcolor='tab:red')

                ax_cv3 = ax_cv.twinx()
                ax_cv3.spines['right'].set_position(('outward', 60))
                ax_cv3.set_ylabel('Power Factor', color='tab:green')
                if f'{turbine_id}_PF' in df_full.columns:
                    ax_cv3.plot(df_full.index, df_full[f'{turbine_id}_PF'], label=f'{turbine_id} Power Factor', color='green', linestyle='-.', alpha=0.7)
                ax_cv3.tick_params(axis='y', labelcolor='tab:green')

                fig_current_voltage.suptitle(f'{turbine_id} Current, Voltage, and Power Factor Over Time')
                fig_current_voltage.legend(loc="upper left", bbox_to_anchor=(0.05,0.95))
                fig_current_voltage.autofmt_xdate()
                plt.grid(True)
                
                img_bytes_cv = fig_to_bytes(fig_current_voltage)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmpfile_cv:
                    tmpfile_cv.write(img_bytes_cv)
                    temp_filename_cv = tmpfile_cv.name
                pdf.image(temp_filename_cv, x=10, w=pdf.w - 20)
                plt.close(fig_current_voltage)
                os.remove(temp_filename_cv)
                pdf.ln(5)

                fig_imbalance, ax_im = plt.subplots(figsize=(15, 7))
                ax_im.plot(df_full.index, df_full[f'{turbine_id}_Imbalance'], label=f'{turbine_id} Current Imbalance', color='teal')
                ax_im.axhline(y=imbalance_threshold, color='red', linestyle=':', label=f'Imbalance Threshold ({imbalance_threshold}%)')
                ax_im.set_xlabel('Date Time')
                ax_im.set_ylabel('Imbalance (%)')
                ax_im.set_title(f'{turbine_id} Current Imbalance Over Time')
                ax_im.legend()
                ax_im.grid(True)
                fig_imbalance.autofmt_xdate()

                img_bytes_im = fig_to_bytes(fig_imbalance)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmpfile_im:
                    tmpfile_im.write(img_bytes_im)
                    temp_filename_im = tmpfile_im.name
                pdf.image(temp_filename_im, x=10, w=pdf.w - 20)
                plt.close(fig_imbalance)
                os.remove(temp_filename_im)
                pdf.ln(10)

            # Save to bytes buffer and return
            output_buffer = BytesIO()
            pdf_bytes = pdf.output(dest='S').encode('latin-1')
            output_buffer.write(pdf_bytes)
            output_buffer.seek(0)
            return output_buffer.getvalue()

        def create_current_ppt_report(analysis_results, df_full, watermark_text, watermark_opacity):
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Advanced Current and Power Analysis in Wind Turbines"
            subtitle.text = "Comprehensive Report"

            if watermark_text:
                left = Inches(1)
                top = Inches(1)
                width = Inches(14)
                height = Inches(7)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = watermark_text
                p.font.size = Pt(80)
                p.font.bold = True
                # Corrección: PP_ALIGN en lugar de PP_RLIGN
                gray_value = int(255 * (1 - watermark_opacity))
                p.font.color.rgb = RGBColor(gray_value, gray_value, gray_value)
                p.alignment = PP_ALIGN.CENTER  # CORREGIDO
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # CORREGIDO

            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "Summary of Analysis Results"

            if watermark_text:
                left = Inches(1)
                top = Inches(1)
                width = Inches(14)
                height = Inches(7)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = watermark_text
                p.font.size = Pt(80)
                p.font.bold = True
                # Corrección: PP_ALIGN en lugar de PP_RLIGN
                gray_value = int(255 * (1 - watermark_opacity))
                p.font.color.rgb = RGBColor(gray_value, gray_value, gray_value)
                p.alignment = PP_ALIGN.CENTER  # CORREGIDO
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # CORREGIDO

            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for turbine_id, data in analysis_results.items():
                p = tf.add_paragraph()
                p.text = f"Turbine {turbine_id}:"
                p.font.size = Pt(18)
                p.font.bold = True

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Peak Total Current: {data['Peak_Total_Current']:.2f} A"
                p.font.size = Pt(16)

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Average Current Imbalance: {data['Avg_Imbalance']:.2f}%"
                p.font.size = Pt(16)

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Average Power Factor: {data['Avg_Power_Factor']:.2f}"
                p.font.size = Pt(16)

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Sustained Overload Events: {len(data['Overload_Events'])}"
                p.font.size = Pt(16)

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Imbalance Events: {len(data['Imbalance_Events'])}"
                p.font.size = Pt(16)

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Low Efficiency Events: {len(data['Efficiency_Events'])}"
                p.font.size = Pt(16)

                p = tf.add_paragraph()
                p.level = 1
                p.text = f"Estimated Total MWh: {data['Annual_Mwh_Estimate']:.2f} MWh"
                p.font.size = Pt(16)

                # Add a slide for detailed events (only Imbalance and Overload)
                if data['Overload_Events'] or data['Imbalance_Events']:
                    event_slide = prs.slides.add_slide(bullet_slide_layout)
                    event_slide.shapes.title.text = f"Detailed Overload & High Imbalance Events for {turbine_id}"
                    
                    if watermark_text:
                        left = Inches(1)
                        top = Inches(1)
                        width = Inches(14)
                        height = Inches(7)
                        txBox = event_slide.shapes.add_textbox(left, top, width, height)
                        tf_wm = txBox.text_frame
                        p_wm = tf_wm.add_paragraph()
                        p_wm.text = watermark_text
                        p_wm.font.size = Pt(80)
                        p_wm.font.bold = True
                        gray_value = int(255 * (1 - watermark_opacity))
                        p_wm.font.color.rgb = RGBColor(gray_value, gray_value, gray_value)
                        p_wm.alignment = PP_ALIGN.CENTER  # CORREGIDO
                        tf_wm.vertical_anchor = MSO_ANCHOR.MIDDLE  # CORREGIDO

                    event_body = event_slide.shapes.placeholders[1]
                    event_tf = event_body.text_frame
                    event_tf.clear()

                    event_data_rows = []
                    for event_type, events in [
                        ("Sustained Overload", data['Overload_Events']),
                        ("High Imbalance", data['Imbalance_Events'])
                    ]:
                        for ts in events:
                            if pd.isna(ts):
                                continue
                            row_data = df_full.loc[ts]
                            event_data_rows.append({
                                "Date_Time": ts.strftime('%Y-%m-%d %H:%M'),
                                "Event_Type": event_type,
                                "Phase_R_Current": f"{row_data.get(f'{turbine_id}A', 0.0):.2f} A",
                                "Phase_S_Current": f"{row_data.get(f'{turbine_id}B', 0.0):.2f} A",
                                "Phase_T_Current": f"{row_data.get(f'{turbine_id}C', 0.0):.2f} A",
                                "Total_Current": f"{row_data.get(f'{turbine_id}_Total', 0.0):.2f} A",
                                "Imbalance": f"{row_data.get(f'{turbine_id}_Imbalance', 0.0):.2f}%",
                                "Power_Factor": f"{row_data.get(f'{turbine_id}_PF', 0.0):.2f}"
                            })
                    
                    # Sort events by Date_Time for chronological order
                    event_data_rows.sort(key=lambda x: datetime.strptime(x["Date_Time"], '%Y-%m-%d %H:%M'))

                    if event_data_rows:
                        # Add table to slide
                        rows = len(event_data_rows) + 1
                        cols = len(event_data_rows[0])
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(15)
                        height = Inches(7)
                        table = event_slide.shapes.add_table(rows, cols, left, top, width, height).table

                        # Set column widths (approximate, adjust as needed)
                        table.columns[0].width = Inches(1.5) # Date_Time
                        table.columns[1].width = Inches(1.5) # Event_Type
                        table.columns[2].width = Inches(1.0) # Phase_R_Current
                        table.columns[3].width = Inches(1.0) # Phase_S_Current
                        table.columns[4].width = Inches(1.0) # Phase_T_Current
                        table.columns[5].width = Inches(1.0) # Total_Current
                        table.columns[6].width = Inches(1.0) # Imbalance
                        table.columns[7].width = Inches(1.0) # Power_Factor

                        # Add headers
                        for i, header in enumerate(event_data_rows[0].keys()):
                            cell = table.cell(0, i)
                            text_frame = cell.text_frame
                            text_frame.text = header.replace("_", " ")
                            text_frame.paragraphs[0].font.size = Pt(9)
                            text_frame.paragraphs[0].font.bold = True
                            # Corrección: PP_ALIGN en lugar de PP_RLIGN
                            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # CORREGIDO

                        # Add data rows
                        for r, row_data in enumerate(event_data_rows):
                            for c, (key, value) in enumerate(row_data.items()):
                                cell = table.cell(r + 1, c)
                                text_frame = cell.text_frame
                                text_frame.text = str(value)
                                text_frame.paragraphs[0].font.size = Pt(8)
                                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # CORREGIDO
                    else:
                        p = event_tf.add_paragraph()
                        p.text = "No 'High Imbalance' or 'Sustained Overload' events detected for this turbine."
                        p.font.size = Pt(16)


            for turbine_id in analysis_results.keys():
                img_slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(img_slide_layout)
                title = slide.shapes.title
                title.text = f"{turbine_id} Current, Voltage, and Power Factor Over Time"

                if watermark_text:
                    left = Inches(1)
                    top = Inches(1)
                    width = Inches(14)
                    height = Inches(7)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = watermark_text
                    p.font.size = Pt(80)
                    p.font.bold = True
                    # Corrección: PP_ALIGN en lugar de PP_RLIGN
                    gray_value = int(255 * (1 - watermark_opacity))
                    p.font.color.rgb = RGBColor(gray_value, gray_value, gray_value)
                    p.alignment = PP_ALIGN.CENTER  # CORREGIDO
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # CORREGIDO

                fig_current_voltage, ax_cv = plt.subplots(figsize=(15, 7))
                ax_cv.set_xlabel('Date Time')
                ax_cv.set_ylabel('Current (A)', color='tab:blue')
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}A'], label=f'{turbine_id} Current R', color='blue', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}B'], label=f'{turbine_id} Current S', color='green', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}C'], label=f'{turbine_id} Current T', color='orange', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}_Total'], label=f'{turbine_id} Total Current', color='purple', linestyle='--', linewidth=2)
                ax_cv.axhline(y=overload_thresholds[turbine_id], color='red', linestyle=':', linewidth=2, label=f'Overload Threshold ({overload_thresholds[turbine_id]}A)')
                ax_cv.tick_params(axis='y', labelcolor='tab:blue')

                ax_cv2 = ax_cv.twinx()
                ax_cv2.set_ylabel('Voltage (V)', color='tab:red')
                if f'{turbine_id}Vab' in df_full.columns:
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vab'], label=f'{turbine_id} Voltage Vab', color='red', alpha=0.5)
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vbc'], label=f'{turbine_id} Voltage Vbc', color='brown', alpha=0.5)
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vca'], label=f'{turbine_id} Voltage Vca', color='pink', alpha=0.5)
                ax_cv2.tick_params(axis='y', labelcolor='tab:red')

                ax_cv3 = ax_cv.twinx()
                ax_cv3.spines['right'].set_position(('outward', 60))
                ax_cv3.set_ylabel('Power Factor', color='tab:green')
                if f'{turbine_id}_PF' in df_full.columns:
                    ax_cv3.plot(df_full.index, df_full[f'{turbine_id}_PF'], label=f'{turbine_id} Power Factor', color='green', linestyle='-.', alpha=0.7)
                ax_cv3.tick_params(axis='y', labelcolor='tab:green')

                fig_current_voltage.suptitle(f'{turbine_id} Current, Voltage, and Power Factor Over Time')
                fig_current_voltage.legend(loc="upper left", bbox_to_anchor=(0.05,0.95))
                fig_current_voltage.autofmt_xdate()
                plt.grid(True)

                img_bytes_cv = fig_to_bytes(fig_current_voltage)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmpfile_cv:
                    tmpfile_cv.write(img_bytes_cv)
                    temp_filename_cv = tmpfile_cv.name
                slide.shapes.add_picture(temp_filename_cv, Inches(1), Inches(2), width=Inches(14), height=Inches(6.5))
                plt.close(fig_current_voltage)
                os.remove(temp_filename_cv)

                slide = prs.slides.add_slide(img_slide_layout)
                title = slide.shapes.title
                title.text = f"{turbine_id} Current Imbalance Over Time"

                if watermark_text:
                    left = Inches(1)
                    top = Inches(1)
                    width = Inches(14)
                    height = Inches(7)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = watermark_text
                    p.font.size = Pt(80)
                    p.font.bold = True
                    # Corrección: PP_ALIGN en lugar de PP_RLIGN
                    gray_value = int(255 * (1 - watermark_opacity))
                    p.font.color.rgb = RGBColor(gray_value, gray_value, gray_value)
                    p.alignment = PP_ALIGN.CENTER  # CORREGIDO
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # CORREGIDO

                fig_imbalance, ax_im = plt.subplots(figsize=(15, 7))
                ax_im.plot(df_full.index, df_full[f'{turbine_id}_Imbalance'], label=f'{turbine_id} Current Imbalance', color='teal')
                ax_im.axhline(y=imbalance_threshold, color='red', linestyle=':', label=f'Imbalance Threshold ({imbalance_threshold}%)')
                ax_im.set_xlabel('Date Time')
                ax_im.set_ylabel('Imbalance (%)')
                ax_im.set_title(f'{turbine_id} Current Imbalance Over Time')
                ax_im.legend()
                ax_im.grid(True)
                fig_imbalance.autofmt_xdate()

                img_bytes_im = fig_to_bytes(fig_imbalance)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmpfile_im:
                    tmpfile_im.write(img_bytes_im)
                    temp_filename_im = tmpfile_im.name
                slide.shapes.add_picture(temp_filename_im, Inches(1), Inches(2), width=Inches(14), height=Inches(6.5))
                plt.close(fig_imbalance)
                os.remove(temp_filename_im)

            output_ppt = BytesIO()
            prs.save(output_ppt)
            output_ppt.seek(0)
            return output_ppt.getvalue()

        def create_current_html_report(analysis_results, df_full, watermark_text, watermark_opacity):
            html_content = """
            <!DOCTYPE html>
            <html>
            <head>
                <title>Advanced Current and Power Analysis Report</title>
                <style>
                    body { font-family: Arial, sans-serif; margin: 20px; position: relative; }
                    .watermark {
                        position: fixed;
                        top: 50%;
                        left: 50%;
                        transform: translate(-50%, -50%) rotate(-45deg);
                        font-size: 100px;
                        color: rgba(192, 192, 192, """ + str(watermark_opacity) + """);
                        z-index: -1;
                        pointer-events: none;
                        white-space: nowrap;
                    }
                    h1 { color: #333; text-align: center; }
                    h2 { color: #555; border-bottom: 2px solid #eee; padding-bottom: 5px; margin-top: 30px; }
                    h3 { color: #777; margin-top: 20px; }
                    .section { margin-bottom: 40px; }
                    .turbine-summary ul { list-style-type: none; padding: 0; }
                    .turbine-summary li { margin-bottom: 5px; }
                    .dataframe { width: 100%; border-collapse: collapse; margin-top: 10px; }
                    .dataframe th, .dataframe td { border: 1px solid #ddd; padding: 8px; text-align: left; }
                    .dataframe th { background-color: #f2f2f2; }
                    .chart-container { margin-top: 20px; text-align: center; }
                    .chart-container img { max-width: 100%; height: auto; border: 1px solid #ddd; }
                </style>
            </head>
            <body>
            """
            if watermark_text:
                html_content += f'<div class="watermark">{watermark_text}</div>'

            html_content += """
                <h1>Advanced Current and Power Analysis Report</h1>

                <div class="section">
                    <h2>1. Summary of Analysis Results</h2>
            """

            for turbine_id, data in analysis_results.items():
                html_content += f"""
                    <div class="turbine-summary">
                        <h3>Turbine {turbine_id} Summary:</h3>
                        <ul>
                            <li>- Peak Current (Phase R): {data['Peak_Current_R']:.2f} A</li>
                            <li>- Peak Current (Phase S): {data['Peak_Current_S']:.2f} A</li>
                            <li>- Peak Current (Phase T): {data['Peak_Current_T']:.2f} A</li>
                            <li>- Peak Total Current: {data['Peak_Total_Current']:.2f} A</li>
                            <li>- Average Current Imbalance: {data['Avg_Imbalance']:.2f}%</li>
                            <li>- Average Power Factor: {data['Avg_Power_Factor']:.2f}</li>
                            <li>- Sustained Overload Events: {len(data['Overload_Events'])}</li>
                            <li>- Imbalance Events: {len(data['Imbalance_Events'])}</li>
                            <li>- Low Efficiency Events: {len(data['Efficiency_Events'])}</li>
                            <li>- Estimated Total MWh over data period: {data['Annual_Mwh_Estimate']:.2f} MWh</li>
                        </ul>
                """
                # --- ONLY High Imbalance and Overload Events for HTML ---
                if data['Overload_Events'] or data['Imbalance_Events']:
                    html_content += f"""
                        <h4>Detailed Overload & High Imbalance Events for {turbine_id}:</h4>
                        <table class="dataframe">
                            <thead>
                                <tr>
                                    <th>Date_Time</th>
                                    <th>Event_Type</th>
                                    <th>Phase_R_Current</th>
                                    <th>Phase_S_Current</th>
                                    <th>Phase_T_Current</th>
                                    <th>Total_Current</th>
                                    <th>Imbalance</th>
                                    <th>Power_Factor</th>
                                </tr>
                            </thead>
                            <tbody>
                    """
                    event_data_rows = []
                    # Include both Overload_Events and Imbalance_Events
                    for event_type, events in [
                        ("Sustained Overload", data['Overload_Events']),
                        ("High Imbalance", data['Imbalance_Events'])
                    ]:
                        for ts in events:
                            if pd.isna(ts):
                                continue
                            row_data = df_full.loc[ts]
                            event_row = [
                                ts.strftime('%Y-%m-%d %H:%M') if not pd.isna(ts) else "Invalid Date",
                                event_type,
                                f"{row_data.get(f'{turbine_id}A', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}B', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}C', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}_Total', 0.0):.2f} A",
                                f"{row_data.get(f'{turbine_id}_Imbalance', 0.0):.2f}%",
                                f"{row_data.get(f'{turbine_id}_PF', 0.0):.2f}"
                            ]
                            event_data_rows.append(event_row)

                    # Sort events by Date_Time for chronological order
                    event_data_rows.sort(key=lambda x: datetime.strptime(x[0], '%Y-%m-%d %H:%M'))

                    for row in event_data_rows:
                        html_content += "<tr>" + "".join([f"<td>{item}</td>" for item in row]) + "</tr>"
                    html_content += """
                            </tbody>
                        </table>
                    """
                else:
                    html_content += "<p>No significant 'High Imbalance' or 'Sustained Overload' events detected in this turbine based on current thresholds.</p>"
                html_content += "</div>"
            html_content += "</div>"

            html_content += """
                <div class="section">
                    <h2>2. Data Visualizations</h2>
            """
            for turbine_id in analysis_results.keys():
                html_content += f"<h3>Graphs for Turbine {turbine_id}:</h3>"

                fig_current_voltage, ax_cv = plt.subplots(figsize=(15, 7))
                ax_cv.set_xlabel('Date Time')
                ax_cv.set_ylabel('Current (A)', color='tab:blue')
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}A'], label=f'{turbine_id} Current R', color='blue', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}B'], label=f'{turbine_id} Current S', color='green', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}C'], label=f'{turbine_id} Current T', color='orange', alpha=0.7)
                ax_cv.plot(df_full.index, df_full[f'{turbine_id}_Total'], label=f'{turbine_id} Total Current', color='purple', linestyle='--', linewidth=2)
                ax_cv.axhline(y=overload_thresholds[turbine_id], color='red', linestyle=':', linewidth=2, label=f'Overload Threshold ({overload_thresholds[turbine_id]}A)')
                ax_cv.tick_params(axis='y', labelcolor='tab:blue')

                ax_cv2 = ax_cv.twinx()
                ax_cv2.set_ylabel('Voltage (V)', color='tab:red')
                if f'{turbine_id}Vab' in df_full.columns:
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vab'], label=f'{turbine_id} Voltage Vab', color='red', alpha=0.5)
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vbc'], label=f'{turbine_id} Voltage Vbc', color='brown', alpha=0.5)
                    ax_cv2.plot(df_full.index, df_full[f'{turbine_id}Vca'], label=f'{turbine_id} Voltage Vca', color='pink', alpha=0.5)
                ax_cv2.tick_params(axis='y', labelcolor='tab:red')

                ax_cv3 = ax_cv.twinx()
                ax_cv3.spines['right'].set_position(('outward', 60))
                ax_cv3.set_ylabel('Power Factor', color='tab:green')
                if f'{turbine_id}_PF' in df_full.columns:
                    ax_cv3.plot(df_full.index, df_full[f'{turbine_id}_PF'], label=f'{turbine_id} Power Factor', color='green', linestyle='-.', alpha=0.7)
                ax_cv3.tick_params(axis='y', labelcolor='tab:green')

                fig_current_voltage.suptitle(f'{turbine_id} Current, Voltage, and Power Factor Over Time')
                fig_current_voltage.legend(loc="upper left", bbox_to_anchor=(0.05,0.95))
                fig_current_voltage.autofmt_xdate()
                plt.grid(True)

                img_bytes_cv = fig_to_bytes(fig_current_voltage)
                img_base64_cv = base64.b64encode(img_bytes_cv).decode('utf-8')
                html_content += f"""
                    <div class="chart-container">
                        <img src="data:image/png;base64,{img_base64_cv}" alt="{turbine_id} Current, Voltage, PF Plot">
                    </div>
                """
                plt.close(fig_current_voltage)

                fig_imbalance, ax_im = plt.subplots(figsize=(15, 7))
                ax_im.plot(df_full.index, df_full[f'{turbine_id}_Imbalance'], label=f'{turbine_id} Current Imbalance', color='teal')
                ax_im.axhline(y=imbalance_threshold, color='red', linestyle=':', label=f'Imbalance Threshold ({imbalance_threshold}%)')
                ax_im.set_xlabel('Date Time')
                ax_im.set_ylabel('Imbalance (%)')
                ax_im.set_title(f'{turbine_id} Current Imbalance Over Time')
                ax_im.legend()
                ax_im.grid(True)
                fig_imbalance.autofmt_xdate()

                img_bytes_im = fig_to_bytes(fig_imbalance)
                img_base64_im = base64.b64encode(img_bytes_im).decode('utf-8')
                html_content += f"""
                    <div class="chart-container">
                        <img src="data:image/png;base64,{img_base64_im}" alt="{turbine_id} Imbalance Plot">
                    </div>
                """
                plt.close(fig_imbalance)

            html_content += """
                </div> </body>
            </html>
            """
            return html_content

        # --- Export Reports ---
        st.header("6. Export Full Analysis Report")
        col1, col2, col3 = st.columns(3)
        with col1:
            st.subheader("PDF Report")
            pdf_bytes = create_current_pdf_report(results, df, watermark_text, watermark_opacity)
            st.download_button(
                label="Download PDF Report",
                data=pdf_bytes,
                file_name="advanced_current_analysis_report.pdf",
                mime="application/pdf"
            )
        with col2:
            st.subheader("PowerPoint Report")
            ppt_bytes = create_current_ppt_report(results, df, watermark_text, watermark_opacity)
            st.download_button(
                label="Download PowerPoint Report",
                data=ppt_bytes,
                file_name="advanced_current_analysis_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        with col3:
            st.subheader("HTML Report")
            html_content = create_current_html_report(results, df, watermark_text, watermark_opacity)
            st.download_button(
                label="Download HTML Report",
                data=html_content,
                file_name="advanced_current_analysis_report.html",
                mime="text/html"
            )
    else:
        st.info("Please upload all required data files to proceed.")
else:
    st.info("Please upload all required data files to start the analysis.")
